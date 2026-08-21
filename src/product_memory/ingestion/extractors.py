from __future__ import annotations

import hashlib
import io
import logging
import re
import zipfile
from collections.abc import Iterator
from dataclasses import dataclass
from datetime import datetime
from email import policy as email_policy
from email.message import EmailMessage
from email.parser import BytesParser
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

import extract_msg
import olefile
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader

from product_memory.ingestion.ocr import OcrEngine

# pypdf logs a WARNING for every recovered xref entry in malformed-but-readable PDFs;
# these are handled automatically and not actionable, so keep them out of app logs.
logging.getLogger("pypdf").setLevel(logging.ERROR)

LOGGER = logging.getLogger(__name__)
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp", ".webp", ".gif"}
SPREADSHEET_ROW_LIMIT = 5000
MAX_EMBEDDED_IMAGE_BYTES = 50 * 1024 * 1024
TABLE_MAX_LABELLED_COLUMNS = 15


@dataclass(slots=True)
class ExtractedDocument:
    content: str
    metadata: dict[str, Any]


class UnreadableDocumentError(ValueError):
    """The file cannot be read at all, for a reason no retry or code change would fix."""


_OLE2_MAGIC = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"
_OOXML_EXTENSIONS = {".docx", ".pptx", ".xlsx"}


def _ooxml_problem(path: Path) -> str | None:
    """Explain why an .docx/.pptx/.xlsx is not the OOXML zip its name promises, if so."""
    with path.open("rb") as handle:
        if handle.read(len(_OLE2_MAGIC)) != _OLE2_MAGIC:
            return None
    try:
        with olefile.OleFileIO(str(path)) as container:
            streams = {"/".join(stream) for stream in container.listdir()}
    except OSError as error:
        return f"unreadable OLE2 container ({error})"
    if any(stream.startswith("EncryptedPackage") for stream in streams):
        return "password protected"
    return "saved in a pre-2007 Office format under a modern extension"


def extract_document(path: Path, ocr: OcrEngine | None = None) -> ExtractedDocument:
    suffix = path.suffix.lower()
    if suffix in _OOXML_EXTENSIONS and (problem := _ooxml_problem(path)) is not None:
        raise UnreadableDocumentError(f"{path}: {problem}")
    collector = _OcrCollector(ocr)
    if suffix == ".pdf":
        extracted = _extract_pdf(path, collector)
    elif suffix == ".docx":
        extracted = _extract_docx(path, collector)
    elif suffix == ".pptx":
        extracted = _extract_pptx(path, collector)
    elif suffix == ".xlsx":
        extracted = _extract_xlsx(path, collector)
    elif suffix == ".msg":
        extracted = _extract_msg(path)
    elif suffix == ".eml":
        extracted = _extract_eml(path, collector)
    elif suffix in IMAGE_EXTENSIONS:
        extracted = _extract_image(path, collector)
    else:
        extracted = ExtractedDocument(content=_read_text(path), metadata={})
    return _sanitize_extracted(extracted)


class _OcrCollector:
    """Reads text from embedded images within a single document, under a fixed budget."""

    def __init__(self, ocr: OcrEngine | None):
        self._ocr = ocr
        self._budget = ocr.settings.ocr_max_images_per_document if ocr else 0
        self._seen_digests: set[str] = set()
        self.sections: list[str] = []
        self.images_seen = 0
        self.images_read = 0
        self.duplicates_skipped = 0

    @property
    def active(self) -> bool:
        return self._ocr is not None and self._ocr.available()

    def add(self, label: str, image: Any) -> None:
        if self._ocr is None:
            return
        self.images_seen += 1
        if self._budget <= 0 or not self._ocr.should_read(image):
            return
        self._budget -= 1
        text = self._ocr.image_to_text(image)
        if text:
            self.images_read += 1
            self.sections.append(f"[Image text: {label}]\n{text}")

    def add_bytes(self, label: str, data: bytes) -> None:
        digest = hashlib.sha256(data).hexdigest()
        if digest in self._seen_digests:
            # Logos and backgrounds repeat on every slide; reading one once is enough.
            self.images_seen += 1
            self.duplicates_skipped += 1
            return
        self._seen_digests.add(digest)
        image = _open_image_bytes(data, label)
        if image is None:
            return
        with image:
            self.add(label, image)

    def compose(self, text: str) -> str:
        return "\n\n".join(part for part in [text, *self.sections] if part)

    def apply_metadata(self, metadata: dict[str, Any]) -> None:
        if not self.active:
            return
        metadata["embedded_image_count"] = self.images_seen
        metadata["ocr_applied"] = bool(self.images_read)
        if self.images_read:
            metadata["ocr_image_count"] = self.images_read
        if self.duplicates_skipped:
            metadata["repeated_image_count"] = self.duplicates_skipped


def _open_image_bytes(data: bytes, label: str) -> Any:
    from PIL import Image

    try:
        image = Image.open(io.BytesIO(data))
        image.load()
    except Exception:
        LOGGER.debug("Skipping unreadable embedded image: %s", label, exc_info=True)
        return None
    return image


def _sanitize_extracted(extracted: ExtractedDocument) -> ExtractedDocument:
    # Postgres text/jsonb columns reject NUL bytes; malformed PDF metadata can contain them.
    return ExtractedDocument(
        content=strip_null_bytes(extracted.content),
        metadata=strip_null_bytes(extracted.metadata),
    )


def strip_null_bytes(value: Any) -> Any:
    if isinstance(value, str):
        return value.replace("\x00", "")
    if isinstance(value, dict):
        return {key: strip_null_bytes(item) for key, item in value.items()}
    if isinstance(value, list):
        return [strip_null_bytes(item) for item in value]
    return value


def _extract_pdf(path: Path, collector: _OcrCollector) -> ExtractedDocument:
    reader = PdfReader(str(path))
    if reader.is_encrypted:
        reader.decrypt("")

    pages = [(page.extract_text() or "").strip() for page in reader.pages]
    text = "\n\n".join(page for page in pages if page)
    metadata: dict[str, Any] = {
        "source_format": "pdf",
        "page_count": len(reader.pages),
    }

    if collector.active:
        for number, page in enumerate(reader.pages, start=1):
            try:
                for image in page.images:
                    if image.image is not None:
                        collector.add(f"page {number}", image.image)
            except Exception:
                LOGGER.debug("Could not read images on page %s of %s", number, path, exc_info=True)
        if not text and collector.images_read:
            metadata["scanned_pdf"] = True

    info = reader.metadata
    if info:
        if info.title:
            metadata["title"] = str(info.title).strip()
        if info.author:
            metadata["author"] = str(info.author).strip()
        if info.subject:
            metadata["subject"] = str(info.subject).strip()
        if info.creator:
            metadata["creator"] = str(info.creator).strip()
        if info.producer:
            metadata["producer"] = str(info.producer).strip()
        if info.creation_date:
            metadata["created_at"] = _isoformat(info.creation_date)
        if info.modification_date:
            metadata["modified_at"] = _isoformat(info.modification_date)

    collector.apply_metadata(metadata)
    return ExtractedDocument(content=collector.compose(text), metadata=metadata)


def _extract_image(path: Path, collector: _OcrCollector) -> ExtractedDocument:
    from PIL import Image

    metadata: dict[str, Any] = {"source_format": "image"}
    with Image.open(path) as image:
        image.load()
        metadata["image_format"] = (image.format or path.suffix.lstrip(".")).lower()
        metadata["image_width"], metadata["image_height"] = image.size
        text = ""
        if collector.active:
            collector.add(path.name, image)
            text = "\n\n".join(collector.sections)
            metadata["ocr_applied"] = bool(collector.images_read)

    return ExtractedDocument(content=text, metadata=metadata)


def _extract_docx(path: Path, collector: _OcrCollector) -> ExtractedDocument:
    document = DocxDocument(str(path))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]

    for table in document.tables:
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        parts.extend(_table_lines(rows))

    metadata: dict[str, Any] = {"source_format": "docx"}
    if collector.active:
        for image_part in document.part.package.image_parts:
            collector.add_bytes(image_part.partname.split("/")[-1], image_part.blob)

    props = document.core_properties
    if props.title:
        metadata["title"] = props.title.strip()
    if props.author:
        metadata["author"] = props.author.strip()
    if props.subject:
        metadata["subject"] = props.subject.strip()
    if props.keywords:
        metadata["keywords"] = props.keywords.strip()
    if props.created:
        metadata["created_at"] = _isoformat(props.created)
    if props.modified:
        metadata["modified_at"] = _isoformat(props.modified)

    collector.apply_metadata(metadata)
    return ExtractedDocument(content=collector.compose("\n\n".join(parts)), metadata=metadata)


def _shape_image_blob(shape: Any) -> bytes | None:
    """Return a shape's embedded image bytes, or None when it has none.

    python-pptx raises ValueError("no embedded image") for a picture whose image is
    linked rather than embedded, so a plain getattr() default is not enough: the
    exception would escape and abort the whole presentation.
    """
    try:
        image = shape.image
    except (AttributeError, ValueError):
        return None
    return None if image is None else image.blob


def _iter_shapes(shapes: Any) -> Iterator[Any]:
    """Yield every shape on a slide, descending into groups.

    Iterating a slide only yields top-level shapes, so anything the author grouped,
    which in practice is most diagrams, would otherwise contribute no text and no image.
    """
    for shape in shapes:
        try:
            is_group = shape.shape_type == MSO_SHAPE_TYPE.GROUP
        except Exception:  # a shape type python-pptx cannot resolve is still worth reading
            is_group = False
        if is_group:
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _extract_pptx(path: Path, collector: _OcrCollector) -> ExtractedDocument:
    presentation = Presentation(str(path))
    parts: list[str] = []

    for number, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)
            if shape.has_table:
                rows = [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                slide_parts.extend(_table_lines(rows))
            if collector.active:
                blob = _shape_image_blob(shape)
                if blob is not None:
                    collector.add_bytes(f"slide {number}", blob)

        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            slide_parts.append(notes)
        if slide_parts:
            parts.append("\n".join(slide_parts))

    metadata: dict[str, Any] = {
        "source_format": "pptx",
        "slide_count": len(presentation.slides),
    }
    props = presentation.core_properties
    if props.title:
        metadata["title"] = props.title.strip()
    if props.author:
        metadata["author"] = props.author.strip()
    if props.subject:
        metadata["subject"] = props.subject.strip()
    if props.keywords:
        metadata["keywords"] = props.keywords.strip()
    if props.created:
        metadata["created_at"] = _isoformat(props.created)
    if props.modified:
        metadata["modified_at"] = _isoformat(props.modified)

    collector.apply_metadata(metadata)
    return ExtractedDocument(content=collector.compose("\n\n".join(parts)), metadata=metadata)


def _extract_xlsx(path: Path, collector: _OcrCollector) -> ExtractedDocument:
    workbook = load_workbook(str(path), data_only=True, read_only=True)
    parts: list[str] = []
    truncated = False
    try:
        for sheet in workbook.worksheets:
            rows: list[list[str]] = []
            for index, row in enumerate(sheet.iter_rows(values_only=True)):
                if index >= SPREADSHEET_ROW_LIMIT:
                    truncated = True
                    break
                cells = ["" if value is None else str(value).strip() for value in row]
                if any(cells):
                    rows.append(cells)
            lines = _table_lines(rows)
            if lines:
                parts.append(f"# {sheet.title}\n" + "\n".join(lines))

        metadata: dict[str, Any] = {
            "source_format": "xlsx",
            "sheet_count": len(workbook.worksheets),
            "sheet_names": [sheet.title for sheet in workbook.worksheets],
        }
        props = workbook.properties
        if props.title:
            metadata["title"] = str(props.title).strip()
        if props.creator:
            metadata["author"] = str(props.creator).strip()
        if props.subject:
            metadata["subject"] = str(props.subject).strip()
        if props.keywords:
            metadata["keywords"] = str(props.keywords).strip()
        if props.created:
            metadata["created_at"] = _isoformat(props.created)
        if props.modified:
            metadata["modified_at"] = _isoformat(props.modified)
    finally:
        workbook.close()

    if truncated:
        metadata["row_limit_reached"] = True

    if collector.active:
        for name, data in _embedded_media(path, "xl/media/"):
            collector.add_bytes(name, data)

    collector.apply_metadata(metadata)
    return ExtractedDocument(content=collector.compose("\n\n".join(parts)), metadata=metadata)


def _table_lines(rows: list[list[str]]) -> list[str]:
    # Chunking splits long tables, so each row repeats its column headers to stay readable alone.
    filled = [row for row in rows if any(cell for cell in row)]
    if not filled:
        return []
    header = filled[0]
    labelled = (
        1 < len(header) <= TABLE_MAX_LABELLED_COLUMNS
        and all(cell for cell in header)
        and len(set(header)) == len(header)
    )
    if not labelled:
        return [" | ".join(cell for cell in row if cell) for row in filled]

    lines = [" | ".join(header)]
    for row in filled[1:]:
        pairs = [
            f"{header[index]}: {cell}"
            for index, cell in enumerate(row)
            if index < len(header) and cell
        ]
        if pairs:
            lines.append(" | ".join(pairs))
    return lines


def _embedded_media(path: Path, prefix: str) -> Iterator[tuple[str, bytes]]:
    with zipfile.ZipFile(path) as archive:
        for info in archive.infolist():
            if not info.filename.startswith(prefix) or info.is_dir():
                continue
            if info.file_size > MAX_EMBEDDED_IMAGE_BYTES:
                LOGGER.debug("Skipping oversized embedded image %s in %s", info.filename, path)
                continue
            yield info.filename.rsplit("/", 1)[-1], archive.read(info)


def _extract_msg(path: Path) -> ExtractedDocument:
    metadata: dict[str, Any] = {"source_format": "msg"}
    with extract_msg.openMsg(str(path)) as message:
        body = (message.body or "").strip()
        if message.subject:
            metadata["title"] = message.subject.strip()
        if message.sender:
            metadata["sender"] = message.sender.strip()
        if message.to:
            metadata["to"] = message.to.strip()
        if message.cc:
            metadata["cc"] = message.cc.strip()
        if message.date:
            metadata["date"] = _isoformat(message.date)
        attachment_names = [name.strip() for name in (a.name for a in message.attachments) if name]
        if attachment_names:
            metadata["attachment_count"] = len(attachment_names)
            metadata["attachment_names"] = attachment_names

    return ExtractedDocument(content=body, metadata=metadata)


def _extract_eml(path: Path, collector: _OcrCollector) -> ExtractedDocument:
    with path.open("rb") as handle:
        message = BytesParser(policy=email_policy.default).parse(handle)

    metadata: dict[str, Any] = {"source_format": "eml"}
    if message["subject"]:
        metadata["title"] = str(message["subject"]).strip()
    for header, key in (("from", "sender"), ("to", "to"), ("cc", "cc")):
        value = message[header]
        if value:
            metadata[key] = str(value).strip()
    sent_at = getattr(message["date"], "datetime", None)
    if sent_at is not None:
        metadata["date"] = _isoformat(sent_at)

    attachment_names: list[str] = []
    for part in message.iter_attachments():
        name = (part.get_filename() or "").strip()
        if name:
            attachment_names.append(name)
        if collector.active and part.get_content_maintype() == "image":
            payload = part.get_payload(decode=True)
            if payload and len(payload) <= MAX_EMBEDDED_IMAGE_BYTES:
                collector.add_bytes(name or "inline image", payload)
    if attachment_names:
        metadata["attachment_count"] = len(attachment_names)
        metadata["attachment_names"] = attachment_names

    collector.apply_metadata(metadata)
    return ExtractedDocument(content=collector.compose(_eml_body(message)), metadata=metadata)


def _eml_body(message: EmailMessage) -> str:
    part = message.get_body(preferencelist=("plain", "html"))
    if part is None:
        return ""
    try:
        text = part.get_content()
    except (LookupError, ValueError):
        payload = part.get_payload(decode=True) or b""
        text = payload.decode("utf-8", errors="replace")
    if part.get_content_subtype() == "html":
        text = html_to_text(text)
    return text.strip()


def html_to_text(html: str) -> str:
    parser = _HtmlTextExtractor()
    parser.feed(html)
    parser.close()
    text = "".join(parser.parts)
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    return re.sub(r"\n\s*\n\s*\n+", "\n\n", text).strip()


class _HtmlTextExtractor(HTMLParser):
    _SKIP = {"script", "style", "head", "title"}
    _BREAK = {"br", "p", "div", "tr", "li", "h1", "h2", "h3", "h4", "h5", "h6", "table", "blockquote"}

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self.parts: list[str] = []
        self._skipping = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        if tag in self._SKIP:
            self._skipping += 1
        elif tag in self._BREAK:
            self.parts.append("\n")
        elif tag == "td":
            self.parts.append(" | ")

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP:
            self._skipping = max(0, self._skipping - 1)
        elif tag in self._BREAK:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if not self._skipping:
            self.parts.append(data)


def _read_text(path: Path) -> str:
    for encoding in ("utf-8", "utf-8-sig", "cp1250", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise UnicodeDecodeError("unknown", b"", 0, 1, f"Cannot decode {path}")


def _isoformat(value: datetime) -> str:
    return value.isoformat()
