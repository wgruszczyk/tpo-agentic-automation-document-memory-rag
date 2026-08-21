from datetime import UTC, datetime
from email.message import EmailMessage
from pathlib import Path

import extract_msg
import olefile
import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from openpyxl.drawing.image import Image as XlsxImage
from PIL import Image
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches

from product_memory.ingestion.extractors import UnreadableDocumentError
from product_memory.ingestion.parser import DocumentParser, EmptyDocumentError
from product_memory.settings import Settings


def test_frontmatter_and_effective_date(tmp_path: Path) -> None:
    path = tmp_path / "meeting.md"
    path.write_text(
        "---\ntitle: Payment meeting\nproject: checkout\neffective_at: 2026-07-31\n---\nBody text",
        encoding="utf-8",
    )
    parser = DocumentParser(Settings(knowledge_dir=tmp_path))
    parsed = parser.parse(path)
    assert parsed.title == "Payment meeting"
    assert parsed.metadata["project"] == "checkout"
    assert parsed.effective_at == datetime(2026, 7, 31, tzinfo=UTC)
    assert parsed.content == "Body text"


def test_date_from_filename(tmp_path: Path) -> None:
    path = tmp_path / "2026-06-15-refinement.txt"
    path.write_text("Useful knowledge", encoding="utf-8")
    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)
    assert parsed.effective_at.date().isoformat() == "2026-06-15"


def test_keeps_text_from_a_note_that_is_only_front_matter(tmp_path: Path) -> None:
    path = tmp_path / "digest.txt"
    path.write_text(
        "---\nWho: Kerstin\nMessage: In TR there are only two brands\n---\n",
        encoding="utf-8",
    )

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert "In TR there are only two brands" in parsed.content
    assert "Who: Kerstin" in parsed.content
    assert parsed.metadata["Message"] == "In TR there are only two brands"


def test_infers_reliable_metadata_from_teams_webvtt(tmp_path: Path) -> None:
    path = tmp_path / "teams-checkout-refinement.vtt"
    path.write_text(
        """WEBVTT

NOTE language:en-US
NOTE duration:"00:12:05.5000000"

Meeting title: Checkout refinement
Date: July 31, 2026 10:00 AM

00:00:00.000 --> 00:00:03.500
<v Alice Brown>We agreed to keep payment retries in the MVP.</v>

00:00:03.500 --> 00:00:05.000
<v Bob Smith>I will update the acceptance criteria.</v>
""",
        encoding="utf-8",
    )

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "Checkout refinement"
    assert parsed.effective_at == datetime(2026, 7, 31, 10, 0, tzinfo=UTC)
    assert parsed.metadata["source_type"] == "meeting_transcript"
    assert parsed.metadata["source_app"] == "microsoft_teams"
    assert parsed.metadata["transcript_format"] == "webvtt"
    assert parsed.metadata["language"] == "en-US"
    assert parsed.metadata["duration_seconds"] == 725.5
    assert parsed.metadata["speakers"] == ["Alice Brown", "Bob Smith"]
    assert parsed.metadata["speaker_count"] == 2


def test_frontmatter_overrides_inferred_transcript_metadata(tmp_path: Path) -> None:
    path = tmp_path / "teams-checkout-refinement.vtt"
    path.write_text(
        """---
title: Explicit title
effective_at: 2026-08-01
project: checkout
---
WEBVTT

Meeting title: Inferred title
Date: July 31, 2026

00:00:00.000 --> 00:00:03.500
<v Alice Brown>We agreed to keep payment retries in the MVP.</v>
""",
        encoding="utf-8",
    )

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "Explicit title"
    assert parsed.effective_at == datetime(2026, 8, 1, tzinfo=UTC)
    assert parsed.metadata["project"] == "checkout"
    assert parsed.metadata["speakers"] == ["Alice Brown"]


def test_ignores_labeled_metadata_inside_transcript_body(tmp_path: Path) -> None:
    path = tmp_path / "2026-06-15-teams-notes.vtt"
    path.write_text(
        """WEBVTT

00:00:00.000 --> 00:00:03.500
<v Alice Brown>Title: this is just something somebody said.</v>

00:00:03.500 --> 00:00:05.000
<v Bob Smith>Date: July 31, 2026 is not the meeting date.</v>
""",
        encoding="utf-8",
    )

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "2026 06 15 teams notes"
    assert parsed.effective_at == datetime(2026, 6, 15, tzinfo=UTC)
    assert parsed.metadata["speakers"] == ["Alice Brown", "Bob Smith"]


def test_source_path_stays_relative_for_symlinked_document(tmp_path: Path) -> None:
    knowledge_dir = tmp_path / "knowledge"
    external_dir = tmp_path / "teams-source"
    external_dir.mkdir(parents=True)
    knowledge_dir.mkdir()
    path = external_dir / "meeting.md"
    path.write_text("Useful knowledge", encoding="utf-8")
    (knowledge_dir / "teams").symlink_to(external_dir, target_is_directory=True)

    parsed = DocumentParser(Settings(knowledge_dir=knowledge_dir)).parse(
        knowledge_dir / "teams" / "meeting.md"
    )

    assert parsed.source_path == "teams/meeting.md"


def test_extracts_docx_content_and_properties(tmp_path: Path) -> None:
    path = tmp_path / "checkout-requirements.docx"
    document = DocxDocument()
    document.core_properties.title = "Checkout Requirements"
    document.core_properties.author = "Product Team"
    document.add_heading("Checkout Requirements", level=1)
    document.add_paragraph("The checkout flow must support saved cards.")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "Priority"
    table.rows[0].cells[1].text = "High"
    document.save(path)

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "Checkout Requirements"
    assert "saved cards" in parsed.content
    assert "Priority | High" in parsed.content
    assert parsed.metadata["author"] == "Product Team"
    assert parsed.metadata["source_format"] == "docx"
    assert parsed.metadata["extension"] == ".docx"


def test_extracts_pdf_content_and_properties(tmp_path: Path) -> None:
    path = tmp_path / "pricing-requirements.pdf"
    _write_simple_pdf(path, "Pricing Requirements", "Discount approval requires finance review.")

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "Pricing Requirements"
    assert "finance review" in parsed.content
    assert parsed.metadata["source_format"] == "pdf"
    assert parsed.metadata["page_count"] == 1
    assert parsed.metadata["extension"] == ".pdf"


def test_extracts_pptx_content_notes_tables_and_properties(tmp_path: Path) -> None:
    path = tmp_path / "checkout-roadmap.pptx"
    presentation = Presentation()
    presentation.core_properties.title = "Checkout Roadmap"
    presentation.core_properties.author = "Product Team"
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Checkout Roadmap"
    slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(1)).text = (
        "Saved cards launch in Q4."
    )
    table = slide.shapes.add_table(1, 2, Inches(1), Inches(3), Inches(5), Inches(1)).table
    table.cell(0, 0).text = "Priority"
    table.cell(0, 1).text = "High"
    slide.notes_slide.notes_text_frame.text = "Confirm launch date with engineering."
    presentation.save(path)

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "Checkout Roadmap"
    assert "Saved cards launch in Q4." in parsed.content
    assert "Priority | High" in parsed.content
    assert "Confirm launch date with engineering." in parsed.content
    assert parsed.metadata["author"] == "Product Team"
    assert parsed.metadata["source_format"] == "pptx"
    assert parsed.metadata["slide_count"] == 1
    assert parsed.metadata["extension"] == ".pptx"


class _FakeAttachment:
    def __init__(self, name: str) -> None:
        self.name = name


class _FakeMsgMessage:
    def __init__(self) -> None:
        self.subject = "Renewal terms for Acme"
        self.sender = "alice@example.com"
        self.to = "bob@example.com"
        self.cc = None
        self.date = datetime(2026, 3, 4, 12, 30, tzinfo=UTC)
        self.body = "We agreed to net-30 payment terms.\n"
        self.attachments = [_FakeAttachment("contract.pdf")]

    def __enter__(self) -> "_FakeMsgMessage":
        return self

    def __exit__(self, *_args: object) -> None:
        return None


def test_extracts_msg_content_and_properties(tmp_path: Path, monkeypatch) -> None:
    path = tmp_path / "renewal.msg"
    path.write_bytes(b"\xd0\xcf\x11\xe0")  # content is irrelevant: the reader is stubbed below
    monkeypatch.setattr(extract_msg, "openMsg", lambda *_args, **_kwargs: _FakeMsgMessage())

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "Renewal terms for Acme"
    assert "net-30 payment terms" in parsed.content
    assert parsed.metadata["source_format"] == "msg"
    assert parsed.metadata["sender"] == "alice@example.com"
    assert parsed.metadata["to"] == "bob@example.com"
    assert parsed.metadata["attachment_names"] == ["contract.pdf"]
    assert parsed.metadata["extension"] == ".msg"
    assert parsed.effective_at == datetime(2026, 3, 4, 12, 30, tzinfo=UTC)


def _write_eml(path: Path, body: str, content_type: str = "text/plain") -> Path:
    subtype = content_type.split("/", 1)[1]
    message = EmailMessage()
    message["Subject"] = "Renewal terms for Acme"
    message["From"] = "alice@example.com"
    message["To"] = "bob@example.com"
    message["Cc"] = "carol@example.com"
    message["Date"] = "Wed, 04 Mar 2026 12:30:00 +0000"
    message.set_content(body, subtype=subtype)
    path.write_bytes(message.as_bytes())
    return path


def test_extracts_eml_content_and_headers(tmp_path: Path) -> None:
    path = _write_eml(tmp_path / "renewal.eml", "We agreed on net-30 payment terms.")

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "Renewal terms for Acme"
    assert "net-30 payment terms" in parsed.content
    assert parsed.metadata["source_format"] == "eml"
    assert parsed.metadata["sender"] == "alice@example.com"
    assert parsed.metadata["to"] == "bob@example.com"
    assert parsed.metadata["cc"] == "carol@example.com"
    assert parsed.metadata["extension"] == ".eml"
    assert parsed.effective_at == datetime(2026, 3, 4, 12, 30, tzinfo=UTC)


def test_reads_an_html_only_email_as_text(tmp_path: Path) -> None:
    body = (
        "<html><head><style>p{color:red}</style></head><body>"
        "<p>Pricing stays at <b>net-30</b>.</p><script>alert(1)</script>"
        "<table><tr><td>DE</td><td>0.36</td></tr></table>"
        "</body></html>"
    )
    path = _write_eml(tmp_path / "html-only.eml", body, content_type="text/html")

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert "Pricing stays at net-30." in parsed.content
    assert "DE | 0.36" in parsed.content
    assert "alert(1)" not in parsed.content
    assert "color:red" not in parsed.content


def test_lists_email_attachments_and_reads_inline_images(tmp_path: Path) -> None:
    picture = _write_png(tmp_path / "chart.png")
    message = EmailMessage()
    message["Subject"] = "Quarterly figures"
    message["Date"] = "Wed, 04 Mar 2026 12:30:00 +0000"
    message.set_content("See the attached chart.")
    message.add_attachment(picture.read_bytes(), maintype="image", subtype="png", filename="chart.png")
    message.add_attachment(b"%PDF-1.4", maintype="application", subtype="pdf", filename="contract.pdf")
    path = tmp_path / "figures.eml"
    path.write_bytes(message.as_bytes())

    parser, ocr = _parser_with_ocr(tmp_path, "Revenue grew by 12 percent")
    parsed = parser.parse(path)

    assert parsed.metadata["attachment_names"] == ["chart.png", "contract.pdf"]
    assert parsed.metadata["attachment_count"] == 2
    assert "Revenue grew by 12 percent" in parsed.content
    assert ocr.calls == 1


class _StubOcr:
    """Stands in for Tesseract so tests stay deterministic and offline."""

    def __init__(self, text: str, settings: Settings) -> None:
        self.text = text
        self.settings = settings
        self.calls = 0

    def available(self) -> bool:
        return True

    def should_read(self, _image: object) -> bool:
        return True

    def image_to_text(self, _image: object) -> str:
        self.calls += 1
        return self.text


def _write_png(path: Path, size: tuple[int, int] = (240, 120)) -> Path:
    Image.new("RGB", size, color="white").save(path)
    return path


def _parser_with_ocr(tmp_path: Path, text: str) -> tuple[DocumentParser, _StubOcr]:
    settings = Settings(knowledge_dir=tmp_path)
    parser = DocumentParser(settings)
    stub = _StubOcr(text, settings)
    parser.ocr = stub
    return parser, stub


def test_extracts_text_from_a_scanned_image(tmp_path: Path) -> None:
    path = _write_png(tmp_path / "invoice-scan.png")
    parser, stub = _parser_with_ocr(tmp_path, "Invoice 2026/03 total 1200 PLN")

    parsed = parser.parse(path)

    assert "Invoice 2026/03 total 1200 PLN" in parsed.content
    assert parsed.metadata["source_format"] == "image"
    assert parsed.metadata["image_width"] == 240
    assert parsed.metadata["image_height"] == 120
    assert parsed.metadata["ocr_applied"] is True
    assert stub.calls == 1


class _MemoryCache:
    def __init__(self) -> None:
        self.entries: dict[str, tuple[str, object]] = {}

    def get(self, source_path: str, signature: str):
        entry = self.entries.get(source_path)
        return entry[1] if entry and entry[0] == signature else None

    def set(self, source_path: str, signature: str, extracted) -> None:
        self.entries[source_path] = (signature, extracted)


def test_unchanged_files_are_not_ocred_again(tmp_path: Path) -> None:
    path = _write_png(tmp_path / "scan.png")
    parser, stub = _parser_with_ocr(tmp_path, "Invoice 2026/03")
    parser.cache = _MemoryCache()

    first = parser.parse(path)
    second = parser.parse(path)

    assert stub.calls == 1
    assert first.content == second.content


def test_rebuild_forces_extraction_to_run_again(tmp_path: Path) -> None:
    path = _write_png(tmp_path / "scan.png")
    parser, stub = _parser_with_ocr(tmp_path, "Invoice 2026/03")
    parser.cache = _MemoryCache()

    parser.parse(path)
    parser.parse(path, force=True)

    assert stub.calls == 2


def test_edited_files_are_extracted_again(tmp_path: Path) -> None:
    path = _write_png(tmp_path / "scan.png")
    parser, stub = _parser_with_ocr(tmp_path, "Invoice 2026/03")
    parser.cache = _MemoryCache()

    parser.parse(path)
    _write_png(path, size=(300, 150))
    parser.parse(path)

    assert stub.calls == 2


def test_image_without_readable_text_is_skipped(tmp_path: Path) -> None:
    path = _write_png(tmp_path / "logo.png")
    parser, _ = _parser_with_ocr(tmp_path, "")

    with pytest.raises(EmptyDocumentError):
        parser.parse(path)


def test_images_are_ignored_when_ocr_is_disabled(tmp_path: Path) -> None:
    path = _write_png(tmp_path / "diagram.png")
    parser = DocumentParser(Settings(knowledge_dir=tmp_path, enable_ocr=False))

    with pytest.raises(EmptyDocumentError):
        parser.parse(path)


def test_reads_text_from_images_embedded_in_a_deck(tmp_path: Path) -> None:
    picture = _write_png(tmp_path / "chart.png")
    path = tmp_path / "quarterly-review.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Review"
    slide.shapes.add_picture(str(picture), Inches(1), Inches(2), Inches(4), Inches(3))
    presentation.save(path)
    parser, stub = _parser_with_ocr(tmp_path, "Revenue grew 18 percent")

    parsed = parser.parse(path)

    assert "Quarterly Review" in parsed.content
    assert "[Image text: slide 1]" in parsed.content
    assert "Revenue grew 18 percent" in parsed.content
    assert parsed.metadata["ocr_applied"] is True
    assert parsed.metadata["ocr_image_count"] == 1
    assert stub.calls == 1


def test_deck_images_are_left_alone_when_ocr_is_disabled(tmp_path: Path) -> None:
    picture = _write_png(tmp_path / "chart.png")
    path = tmp_path / "quarterly-review.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Quarterly Review"
    slide.shapes.add_picture(str(picture), Inches(1), Inches(2), Inches(4), Inches(3))
    presentation.save(path)

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path, enable_ocr=False)).parse(path)

    assert parsed.content.strip() == "Quarterly Review"
    assert "ocr_applied" not in parsed.metadata


def test_deck_with_a_linked_image_is_still_parsed(tmp_path: Path) -> None:
    picture_file = _write_png(tmp_path / "chart.png")
    path = tmp_path / "linked-image.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Linked Image Deck"
    picture = slide.shapes.add_picture(str(picture_file), Inches(1), Inches(2), Inches(4), Inches(3))
    # A linked rather than embedded picture carries no r:embed, and python-pptx then
    # raises ValueError("no embedded image") when .image is read.
    del picture._element.blipFill.blip.attrib[qn("r:embed")]
    presentation.save(path)
    parser, stub = _parser_with_ocr(tmp_path, "Revenue grew 18 percent")

    parsed = parser.parse(path)

    assert "Linked Image Deck" in parsed.content
    assert stub.calls == 0


def test_ocr_stops_after_the_image_budget_is_reached(tmp_path: Path) -> None:
    picture = _write_png(tmp_path / "chart.png")
    path = tmp_path / "many-images.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Many Images"
    for offset in range(3):
        slide.shapes.add_picture(str(picture), Inches(1), Inches(1 + offset), Inches(2), Inches(1))
    presentation.save(path)

    settings = Settings(knowledge_dir=tmp_path, ocr_max_images_per_document=2)
    parser = DocumentParser(settings)
    stub = _StubOcr("Revenue grew 18 percent", settings)
    parser.ocr = stub

    parsed = parser.parse(path)

    assert stub.calls == 2
    assert parsed.metadata["embedded_image_count"] == 3
    assert parsed.metadata["ocr_image_count"] == 2


def test_extracts_spreadsheet_cells_sheets_and_properties(tmp_path: Path) -> None:
    path = tmp_path / "commercial-offer.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Pricing"
    sheet.append(["Service", "Fee"])
    sheet.append(["Transaction", 0.35])
    second = workbook.create_sheet("Notes")
    second.append(["Valid until 2026-12-31"])
    workbook.properties.title = "Commercial Offer"
    workbook.properties.creator = "Finance"
    workbook.save(path)

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert parsed.title == "Commercial Offer"
    assert "Service | Fee" in parsed.content
    assert "Service: Transaction | Fee: 0.35" in parsed.content
    assert "Valid until 2026-12-31" in parsed.content
    assert parsed.metadata["source_format"] == "xlsx"
    assert parsed.metadata["sheet_count"] == 2
    assert parsed.metadata["sheet_names"] == ["Pricing", "Notes"]
    assert parsed.metadata["author"] == "Finance"
    assert parsed.metadata["extension"] == ".xlsx"


def test_spreadsheet_rows_repeat_their_column_headers(tmp_path: Path) -> None:
    path = tmp_path / "fees.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Payment fees EU"
    sheet.append(["Country", "Method", "Fee"])
    for number in range(40):
        sheet.append([f"Country {number}", "Card", f"{number}.10 EUR"])
    workbook.save(path)

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    # A late row must still name its columns, because chunking cuts the sheet away from its header.
    assert "Country: Country 39 | Method: Card | Fee: 39.10 EUR" in parsed.content


def test_spreadsheet_without_a_usable_header_keeps_plain_rows(tmp_path: Path) -> None:
    path = tmp_path / "notes.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["Country", None, "Fee"])
    sheet.append(["Poland", "Card", "1.10 EUR"])
    workbook.save(path)

    parsed = DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

    assert "Country | Fee" in parsed.content
    assert "Poland | Card | 1.10 EUR" in parsed.content


def test_reads_text_from_images_embedded_in_a_spreadsheet(tmp_path: Path) -> None:
    picture = _write_png(tmp_path / "chart.png")
    path = tmp_path / "with-chart.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    sheet.append(["Quarterly figures"])
    sheet.add_image(XlsxImage(str(picture)), "C3")
    workbook.save(path)
    parser, stub = _parser_with_ocr(tmp_path, "Revenue grew 18 percent")

    parsed = parser.parse(path)

    assert "Quarterly figures" in parsed.content
    assert "Revenue grew 18 percent" in parsed.content
    assert parsed.metadata["ocr_image_count"] == 1
    assert stub.calls == 1


def _write_simple_pdf(path: Path, title: str, body: str) -> None:
    lines = [title, body]
    text_commands = ["BT", "/F1 12 Tf", "72 720 Td"]
    for index, line in enumerate(lines):
        if index:
            text_commands.append("0 -18 Td")
        text_commands.append(f"({_escape_pdf_text(line)}) Tj")
    text_commands.append("ET")
    stream = "\n".join(text_commands).encode("ascii")

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Title (" + _escape_pdf_text(title).encode("ascii") + b") >>",
    ]

    content = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for number, pdf_object in enumerate(objects, start=1):
        offsets.append(len(content))
        content.extend(f"{number} 0 obj\n".encode("ascii"))
        content.extend(pdf_object)
        content.extend(b"\nendobj\n")

    xref_offset = len(content)
    content.extend(f"xref\n0 {len(objects) + 1}\n".encode("ascii"))
    content.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        content.extend(f"{offset:010d} 00000 n \n".encode("ascii"))
    content.extend(
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R /Info 6 0 R >>\n"
        f"startxref\n{xref_offset}\n%%EOF\n".encode("ascii")
    )
    path.write_bytes(bytes(content))


def _escape_pdf_text(value: str) -> str:
    return value.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")


def test_empty_file_is_skipped_rather_than_failing_the_reader(tmp_path: Path) -> None:
    path = tmp_path / "placeholder.docx"
    path.write_bytes(b"")

    with pytest.raises(EmptyDocumentError):
        DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)


class _FakeOleContainer:
    def __init__(self, streams: list[list[str]]) -> None:
        self._streams = streams

    def __enter__(self) -> "_FakeOleContainer":
        return self

    def __exit__(self, *_args: object) -> None:
        return None

    def listdir(self) -> list[list[str]]:
        return self._streams


def _write_ole2_stub(path: Path) -> None:
    path.write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 64)


def test_password_protected_workbook_is_skipped_with_a_reason(tmp_path: Path, monkeypatch) -> None:
    # Office encrypts an OOXML file by wrapping it in an OLE2 container, so it is no longer a zip.
    path = tmp_path / "test users.xlsx"
    _write_ole2_stub(path)
    monkeypatch.setattr(
        olefile, "OleFileIO", lambda *_a, **_k: _FakeOleContainer([["EncryptedPackage"], ["EncryptionInfo"]])
    )

    with pytest.raises(UnreadableDocumentError, match="password protected"):
        DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)


def test_pre_2007_workbook_named_xlsx_is_skipped_with_a_reason(tmp_path: Path, monkeypatch) -> None:
    path = tmp_path / "old numbers.xlsx"
    _write_ole2_stub(path)
    monkeypatch.setattr(olefile, "OleFileIO", lambda *_a, **_k: _FakeOleContainer([["Workbook"]]))

    with pytest.raises(UnreadableDocumentError, match="pre-2007"):
        DocumentParser(Settings(knowledge_dir=tmp_path)).parse(path)

